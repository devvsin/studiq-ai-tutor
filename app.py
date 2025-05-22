import os
import time
import uuid
import json
import logging
import re
from flask import Flask, render_template, request, redirect, url_for, flash, session, jsonify, send_from_directory
from flask_cors import CORS
from pymongo import MongoClient
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

# Try to import Gemini AI (Google's API)
try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False
    logging.warning("Google Generative AI module not available. Advanced AI features will be disabled.")

# Try to import optional modules
try:
    from gtts import gTTS
    GTTS_AVAILABLE = True
except ImportError:
    GTTS_AVAILABLE = False
    logging.warning("gTTS not available. Audio generation will be disabled.")

# Try to import PDF reader
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logging.warning("PyPDF2 not available. PDF processing will be disabled.")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv()

# Initialize Flask application
app = Flask(__name__, static_folder='static', static_url_path='/static')
CORS(app)  # Enable CORS for API calls
app.secret_key = os.getenv("SECRET_KEY", "study_q_app_secret_key_2025")

# MongoDB setup
mongo_uri = os.getenv('MONGO_URI', 'your mongo uri')
client = MongoClient(mongo_uri)
db = client['hack']  # Database name
users_collection = db['users']  # Collection for users
quiz_collection = db['quiz']  # Collection for quiz results

# Configure directories
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Audio directory for TTS
AUDIO_FOLDER = os.path.join('static', 'audio')
if not os.path.exists(AUDIO_FOLDER):
    os.makedirs(AUDIO_FOLDER)

# Maximum allowed file size (20MB)
app.config['MAX_CONTENT_LENGTH'] = 20 * 1024 * 1024

# Configure Gemini API if available
if HAS_GEMINI:
    try:
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel("gemini-1.5-flash")
        logging.info("Gemini AI configured successfully")
    except Exception as e:
        logger.error(f"Failed to configure Gemini API: {str(e)}")
        HAS_GEMINI = False

# Simple in-memory storage for user sessions and documents
user_sessions = {}
document_content = {}

# Learning styles with detailed prompts
LEARNING_STYLES = {
    "visual": {
        "prompt": "Use visual descriptions, diagrams, and spatial organization to explain concepts. Organize information with clear headings, bullet points, and visual markers. Describe concepts in terms of how they would look visually.",
        "description": "Visual learner who prefers diagrams, images, and spatial information",
        "example": "Imagine a flowchart showing...",
        "model_instruction": "Create visually structured explanations with clear headings, bullet points, emojis as visual markers, and spatial organization to help with mental mapping."
    },
    "auditory": {
        "prompt": "Explain clearly in a conversational tone as if speaking aloud to the student. Use rhythm, repetition, clear transitions, and storytelling approaches. Focus on how ideas sound when explained verbally.",
        "description": "Auditory learner who processes information best through hearing and speaking",
        "example": "Listen to how these concepts connect...",
        "model_instruction": "use a conversational tone with clear transitions and storytelling techniques"
    },
    "hands-on": {
        "prompt": "Provide actionable exercises, practical examples, and step-by-step instructions. Break down concepts into clear steps that can be practiced, focusing on how the student can apply the knowledge immediately.",
        "description": "Kinesthetic learner who prefers learning by doing and practical examples",
        "example": "Try this exercise to understand the concept...",
        "model_instruction": "include practical examples, step-by-step instructions, and actionable exercises"
    },
    "reading": {
        "prompt": "Provide detailed, well-structured written explanations with logical flow. Include relevant context, nuanced details, and proper citations when applicable. Focus on precise language and thorough coverage of concepts.",
        "description": "Reader who prefers comprehensive written explanations",
        "example": "The following elements are crucial to understand...",
        "model_instruction": "write detailed, well-structured explanations with logical flow and precise language"
    },
    "blended": {
        "prompt": "Combine visual elements, clear explanations, practical examples, and well-structured text. Provide a balanced mix of learning approaches, addressing different aspects of learning. Include both theoretical explanations and practical applications.",
        "description": "Balanced learner who benefits from mixed learning approaches",
        "example": "Consider both the theory and practice of this concept...",
        "model_instruction": "provide a balanced mix of visual elements, clear explanations, and practical examples"
    }
}

# Study environments and focus preferences
STUDY_ENVIRONMENTS = {
    "Quiet room": "I'll structure my explanations for deep focus, with clear breaks and sections to help you concentrate in your quiet study space.",
    "With music": "I'll structure explanations into rhythm-friendly chunks that work well when studying with background music.",
    "Café": "I'll provide explanations that work well in busy environments with potential distractions - using clear highlights and summaries.",
    "Library": "I'll provide comprehensive, well-structured content optimized for deep, focused study sessions in quiet environments."
}

FOCUS_DURATIONS = {
    "less than 15 minutes": "I'll keep explanations very concise with key points highlighted at the beginning.",
    "15-30 minutes": "I'll provide medium-length explanations with clear structure and breaks.",
    "30-60 minutes": "I'll offer more comprehensive explanations with multiple examples and deeper connections.",
    "More than an hour": "I'll provide in-depth content with multiple perspectives, examples, and connections to other topics."
}

HELP_APPROACHES = {
    "Step-by-step guidance": "When you're stuck, I'll guide you through solutions one clear step at a time, checking understanding at each stage.",
    "A detailed explanation": "When you're stuck, I'll provide thorough explanations with context and principles to deepen your understanding.",
    "Hints or clues": "When you're stuck, I'll offer subtle hints that lead you to discover the solution yourself.",
    "A summary or analogy": "When you face challenges, I'll use analogies and simplified summaries to make concepts more intuitive.",
    "Re-explanation from a new angle": "If you're struggling, I'll approach the topic from completely different perspectives until we find what clicks for you."
}

INTERACTION_STYLES = {
    "Friendly and casual": "I'll maintain a warm, conversational tone with encouraging language.",
    "Professional and straightforward": "I'll keep explanations clear, concise and focused, maintaining a professional tone.",
    "Encouraging and motivational": "I'll emphasize progress, potential, and positive reinforcement throughout our interactions.",
    "Witty and humorous": "I'll incorporate appropriate humor and interesting facts to keep learning engaging."
}

# Helper functions

# Email validation function
def is_valid_email(email):
    pattern = r'^[\w\.-]+@[\w\.-]+\.\w+$'
    return re.match(pattern, email) is not None

# Get or create a session ID
def get_session_id():
    """Get current session ID or create a new one"""
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
        init_new_session(session['session_id'])
    
    session_id = session['session_id']
    
    # Ensure the session exists in our storage
    if session_id not in user_sessions:
        init_new_session(session_id)
    
    return session_id

# Initialize a new session with welcome message
def init_new_session(session_id):
    """Initialize a new user session with default welcome message"""
    welcome_msg = "Hello! I'm your StudiQ AI tutor. How can I help you with your studies today?"
    
    # Get learning style if available from quiz data
    learning_style = "blended"  # Default
    
    if 'user_email' in session:
        quiz_data = quiz_collection.find_one({"user_email": session['user_email']})
        
        if quiz_data:
            # Map quiz responses to learning styles
            style_mapping = {
                "Watching videos": "visual",
                "Reading books": "reading",
                "Listening to podcasts": "auditory",
                "Doing it myself": "hands-on"
            }
            
            if 'learningStyle' in quiz_data:
                learning_style = style_mapping.get(quiz_data['learningStyle'], "blended")
    
    user_sessions[session_id] = {
        'learning_style': learning_style,
        'documents': [],
        'chat_history': [
            {
                'role': 'assistant',
                'content': welcome_msg,
                'timestamp': time.time()
            }
        ],
        'last_active': time.time()
    }
    return user_sessions[session_id]

# Check if user has completed the quiz
def has_completed_quiz(email):
    """Check if a user has completed the quiz"""
    if not email:
        return False
    quiz_result = quiz_collection.find_one({"user_email": email})
    return quiz_result is not None

# File processing functions
def extract_text_from_pdf(pdf_path):
    """Extract text from PDF file"""
    if not PDF_AVAILABLE:
        return "PDF processing is not available."
        
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}")
        return ""

def process_file(file_path):
    """Process uploaded file to extract text"""
    if not os.path.exists(file_path):
        return "File not found"
        
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.txt', '.md', '.csv']:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            return ""
    # Basic PowerPoint support
    elif ext in ['.ppt', '.pptx']:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"\nSlide {i+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error reading PowerPoint file: {str(e)}")
            return f"Error processing PowerPoint: {str(e)}"
    else:
        return f"Unsupported file format: {ext}"

# Audio generation function
def generate_audio_from_text(text, session_id=None):
    """Generate audio file from text using gTTS"""
    if not GTTS_AVAILABLE:
        return None
        
    try:
        # Clean text for audio (remove markdown, etc.)
        clean_text = text
        clean_text = re.sub(r'\*\*(.*?)\*\*', r'\1', clean_text)  # Remove bold
        clean_text = re.sub(r'\*(.*?)\*', r'\1', clean_text)       # Remove italic
        clean_text = re.sub(r'#{1,6}\s+(.*?)\n', r'\1. ', clean_text)  # Convert headers to sentences
        
        # Generate unique filename
        filename = f"{session_id or uuid.uuid4()}_audio_{int(time.time())}.mp3"
        file_path = os.path.join(AUDIO_FOLDER, filename)
        
        # Generate audio
        tts = gTTS(text=clean_text, lang='en', slow=False)
        tts.save(file_path)
        
        # Return URL for the audio file
        audio_url = f"/audio/{filename}"
        
        return {
            "url": audio_url,
            "file_path": file_path
        }
    except Exception as e:
        logger.error(f"Error generating audio: {str(e)}")
        return None

# Visual enhancements for visual learners
def add_visual_elements(text):
    """Add visual elements to the response for visual learners"""
    # Add emoji indicators to lists
    text = re.sub(r'^- ', r'📌 ', text, flags=re.MULTILINE)
    
    # Add emojis to headings based on content
    def add_heading_emoji(match):
        heading = match.group(2)
        heading_lower = heading.lower()
        
        emoji = "🔍"  # Default
        
        # Map common topics to relevant emojis
        if any(word in heading_lower for word in ["math", "equation", "calculation", "number"]):
            emoji = "🧮"
        elif any(word in heading_lower for word in ["history", "past", "ancient", "timeline"]):
            emoji = "📜"
        elif any(word in heading_lower for word in ["science", "physics", "chemistry", "biology"]):
            emoji = "🔬"
        elif any(word in heading_lower for word in ["step", "how to", "process", "procedure"]):
            emoji = "📋"
        elif any(word in heading_lower for word in ["example", "sample", "instance"]):
            emoji = "💡"
        elif any(word in heading_lower for word in ["summary", "conclusion", "recap"]):
            emoji = "📝"
            
        return f"{match.group(1)} {emoji} {heading}"
    
    text = re.sub(r'^(#{1,3}) (.+)$', add_heading_emoji, text, flags=re.MULTILINE)
    
    # Add dividers between sections
    text = re.sub(r'\n#{2,3} ', r'\n\n---\n\n## ', text)
    
    return text

def generate_visual_prompt(base_prompt, learning_style="visual"):
    """Enhance the prompt for visual learners"""
    visual_addition = """
For visual learners, focus on:
- Using emojis as visual indicators
- Creating clear hierarchical structure with headings
- Using bullet points and numbered lists for steps
- Creating visual separation between concepts
- Using spatial organization to illustrate relationships
    """
    return base_prompt + visual_addition

# Generate personalized learning prompt based on quiz answers
def generate_personalized_prompt(quiz_data):
    """Generate a personalized learning prompt based on quiz answers"""
    if not quiz_data:
        return "I'll adapt to your learning style as we interact."
    
    personalization = []
    
    # Learning style
    if 'learningStyle' in quiz_data:
        style_mapping = {
            "Watching videos": "visual",
            "Reading books": "reading",
            "Listening to podcasts": "auditory",
            "Doing it myself": "hands-on"
        }
        
        learning_style = style_mapping.get(quiz_data['learningStyle'], "blended")
        personalization.append(f"Your primary learning style appears to be: {learning_style}.")
    
    # Study environment
    if 'studyEnv' in quiz_data:
        env_guidance = STUDY_ENVIRONMENTS.get(quiz_data['studyEnv'], "")
        if env_guidance:
            personalization.append(env_guidance)
    
    # Information retention
    if 'retainInfo' in quiz_data:
        if quiz_data['retainInfo'] == "Writing notes":
            personalization.append("I'll structure content in a way that's easy to take notes on, with clear key points and organized sections.")
        elif quiz_data['retainInfo'] == "Applying it in practice":
            personalization.append("I'll emphasize practical applications and provide examples you can try yourself.")
        elif quiz_data['retainInfo'] == "Teaching others":
            personalization.append("I'll frame explanations in a way that prepares you to explain concepts to others.")
        elif quiz_data['retainInfo'] == "Rewatching videos":
            personalization.append("I'll create explanations with clear visual markers and memorable structure.")
    
    # Topic approach
    if 'newTopicPreference' in quiz_data:
        if "broad overview" in quiz_data['newTopicPreference'].lower():
            personalization.append("I'll start with big-picture concepts before diving into details.")
        elif "dive into details" in quiz_data['newTopicPreference'].lower():
            personalization.append("I'll focus on specific details and mechanisms right away.")
    
    # Study duration
    if 'studyDuration' in quiz_data:
        duration_guidance = FOCUS_DURATIONS.get(quiz_data['studyDuration'], "")
        if duration_guidance:
            personalization.append(duration_guidance)
    
    # Help approach
    if 'struggleHelp' in quiz_data:
        help_guidance = HELP_APPROACHES.get(quiz_data['struggleHelp'], "")
        if help_guidance:
            personalization.append(help_guidance)
    
    # Interaction style
    if 'chatbotInteraction' in quiz_data:
        interaction_style = INTERACTION_STYLES.get(quiz_data['chatbotInteraction'], "")
        if interaction_style:
            personalization.append(interaction_style)
    
    # Combine all personalizations
    if personalization:
        return " ".join(personalization)
    else:
        return "I'll adapt to your learning style as we interact."

# Make functions available to templates
@app.context_processor
def utility_processor():
    return {
        'has_completed_quiz': has_completed_quiz
    }

# Routes for the application

# Home page route
@app.route('/')
def home():
    return render_template('landland.html')

# Login page route
@app.route('/loginpage')
def loginpage():
    return render_template('login2.html')

# Login form submission
@app.route('/login', methods=['POST'])
def login():
    try:
        email = request.form.get('login-email')
        password = request.form.get('login-password')

        if not all([email, password]):
            flash('Please enter both email and password.', 'error')
            return redirect(url_for('loginpage'))

        # Check if user exists and password matches
        user = users_collection.find_one({"email": email})

        if user and check_password_hash(user["password"], password):
            # Store user info in session
            session['user_email'] = email
            session['user_role'] = user['role']
            session['user_name'] = user['full_name']
            
            # Check if user has completed the quiz
            quiz_result = quiz_collection.find_one({"user_email": email})
            
            if quiz_result:
                # User has completed the quiz, redirect to main app
                return redirect(url_for('app_dashboard'))
            else:
                # User needs to take the quiz
                return redirect(url_for('quiz'))
        else:
            flash('Invalid email or password.', 'error')
            return redirect(url_for('loginpage'))

    except Exception as e:
        logger.error(f"Login error: {str(e)}")
        flash('An error occurred during login. Please try again.', 'error')
        return redirect(url_for('loginpage'))

# Signup form submission
@app.route('/signup', methods=['POST'])
def signup():
    try:
        full_name = request.form.get('signup-name')
        email = request.form.get('signup-email')
        password = request.form.get('signup-password')
        role = request.form.get('student-role', 'student')  # Get the selected role

        # Validation
        if not all([full_name, email, password]):
            flash('All fields are required.', 'error')
            return redirect(url_for('loginpage'))

        if not is_valid_email(email):
            flash('Please enter a valid email address.', 'error')
            return redirect(url_for('loginpage'))

        if len(password) < 6:
            flash('Password must be at least 6 characters long.', 'error')
            return redirect(url_for('loginpage'))

        # Check if user already exists
        if users_collection.find_one({"email": email}):
            flash('Email already exists. Please sign in.', 'error')
            return redirect(url_for('loginpage'))

        # Hash the password before saving it
        hashed_password = generate_password_hash(password)

        # Create a new user document
        user = {
            "full_name": full_name,
            "email": email,
            "password": hashed_password,
            "role": role
        }

        # Insert the user into the database
        result = users_collection.insert_one(user)

        if result.inserted_id:
            # Store user info in session
            session['user_email'] = email
            session['user_role'] = role
            session['user_name'] = full_name

            # Redirect to quiz page after account creation
            flash('Account created successfully! Welcome to StudiQ.', 'success')
            return redirect(url_for('quiz'))
        else:
            flash('Error creating account. Please try again.', 'error')
            return redirect(url_for('loginpage'))

    except Exception as e:
        logger.error(f"Signup error: {str(e)}")
        flash('An error occurred during signup. Please try again.', 'error')
        return redirect(url_for('loginpage'))

# Landing page route (after login)
@app.route('/landland')
def landland():
    return render_template('landland.html')

# Academics route - redirects based on user status
@app.route('/academics')
def academics():
    if 'user_email' not in session:
        # Not logged in, redirect to login
        return redirect(url_for('loginpage'))
    
    # Check if quiz is completed
    if has_completed_quiz(session['user_email']):
        # Quiz completed, go to app dashboard
        return redirect(url_for('app_dashboard'))
    else:
        # Quiz not completed, go to quiz
        return redirect(url_for('quiz'))

# Quiz page route
@app.route('/quiz')
def quiz():
    if 'user_email' not in session:
        flash('Please login first.', 'error')
        return redirect(url_for('loginpage'))
    return render_template('quiz2.html')

# Save quiz results
@app.route('/save_quiz_results', methods=['POST'])
def save_quiz_results():
    if 'user_email' not in session:
        return jsonify({'success': False, 'error': 'Not logged in'}), 401

    data = request.get_json()
    if not data:
        return jsonify({'success': False, 'error': 'No data received'}), 400

    # Build the quiz document with all answers as separate fields
    quiz_doc = {'user_email': session['user_email']}
    for key, value in data.items():
        quiz_doc[key] = value

    # Store or update quiz results
    quiz_collection.update_one(
        {'user_email': session['user_email']}, 
        {'$set': quiz_doc}, 
        upsert=True
    )

    return jsonify({'success': True})

# Main app dashboard
@app.route('/app')
def app_dashboard():
    if 'user_email' not in session:
        flash('Please login first.', 'error')
        return redirect(url_for('loginpage'))
    
    # Check if quiz is completed
    if not has_completed_quiz(session['user_email']):
        # Quiz not completed, redirect to quiz
        flash('Please complete the learning style assessment first.', 'info')
        return redirect(url_for('quiz'))
        
    return render_template('index.html')

# Logout route
@app.route('/logout')
def logout():
    session.clear()
    flash('You have been logged out.', 'success')
    return redirect(url_for('home'))

# Static files route
@app.route('/static/<path:path>')
def serve_static(path):
    return send_from_directory('static', path)

# Audio files route
@app.route('/audio/<path:filename>')
def serve_audio(filename):
    return send_from_directory(AUDIO_FOLDER, filename)

# API Routes

# Health check endpoint
@app.route('/api/healthcheck', methods=['GET'])
def healthcheck():
    """Basic health check endpoint"""
    return jsonify({
        "status": "ok", 
        "timestamp": time.time(),
        "gemini_available": HAS_GEMINI,
        "tts_available": GTTS_AVAILABLE
    })

# File upload endpoint
@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Upload and process a document"""
    # Check if user is logged in
    if 'user_email' not in session:
        return jsonify({"error": "You must be logged in to upload files"}), 401
        
    try:
        session_id = get_session_id()
        
        # Check if file is present
        if 'file' not in request.files:
            return jsonify({"error": "No file part"}), 400
            
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No selected file"}), 400
            
        # Secure the filename
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], f"{session_id}_{filename}")
        
        # Save the file
        file.save(filepath)
        logger.info(f"File saved: {filepath}")
        
        # Process file to extract text
        text = process_file(filepath)
        
        if not text:
            # Delete the file if we couldn't extract text
            if os.path.exists(filepath):
                os.remove(filepath)
            return jsonify({"error": "Failed to extract text from file"}), 400
        
        # Store document info
        doc_id = str(uuid.uuid4())
        document_content[doc_id] = text
        
        # Add to user session
        user_sessions[session_id]['documents'].append({
            'id': doc_id,
            'filename': filename,
            'filepath': filepath,
            'upload_time': time.time()
        })
        
        # Add system message about document upload
        user_sessions[session_id]['chat_history'].append({
            'role': 'assistant',
            'content': f"I've processed your document \"{filename}\". You can now ask me questions about it!",
            'timestamp': time.time()
        })
        
        return jsonify({
            "success": True,
            "document_id": doc_id,
            "filename": filename,
            "text_length": len(text)
        })
        
    except Exception as e:
        logger.error(f"Error in upload_file: {str(e)}")
        return jsonify({"error": str(e)}), 500

# Chat API endpoint
@app.route('/api/chat', methods=['POST'])
def chat():
    """Handle chat messages with improved formatting and learning style adaptations"""
    try:
        session_id = get_session_id()
        session_data = user_sessions[session_id]
        
        # Parse request data
        data = request.json
        if not data or 'message' not in data:
            return jsonify({"error": "No message provided"}), 400
            
        user_message = data['message']
        learning_style = session_data.get('learning_style', 'blended')
        
        # Add user message to history
        session_data['chat_history'].append({
            'role': 'user',
            'content': user_message,
            'timestamp': time.time()
        })
        
        # If Gemini is not available, return a fallback response
        if not HAS_GEMINI:
            fallback_response = "I'm sorry, but the advanced AI model is not available right now. Please check the API key configuration or try again later."
            session_data['chat_history'].append({
                'role': 'assistant',
                'content': fallback_response,
                'timestamp': time.time()
            })
            return jsonify({
                "response": fallback_response,
                "timestamp": time.time()
            })
        
        # Prepare context from documents
        context = ""
        has_document = False
        if session_data['documents']:
            # Use the most recent document as context
            doc_id = session_data['documents'][-1]['id']
            if doc_id in document_content:
                context = document_content[doc_id][:5000]  # Limit context size
                has_document = True
        
        # Prepare chat history for context
        chat_context = ""
        for msg in session_data['chat_history'][-5:]:  # Use last 5 messages
            if msg['role'] == 'user':
                chat_context += f"User: {msg['content']}\n"
            else:
                chat_context += f"AI: {msg['content']}\n"
        
        # Get user quiz data for personalization
        personalized_instruction = "I'll adapt to your learning style as we interact."
        if 'user_email' in session:
            quiz_data = quiz_collection.find_one({"user_email": session['user_email']})
            if quiz_data:
                personalized_instruction = generate_personalized_prompt(quiz_data)
        
        # Prepare prompt for Gemini with formatting instructions
        # Get learning style specific instructions
        style_info = LEARNING_STYLES.get(learning_style, LEARNING_STYLES['blended'])
        learning_instruction = style_info["prompt"]
        model_instruction = style_info["model_instruction"]
        
        # Base prompt
        base_prompt = f"""
You are a helpful AI tutor assistant named StudiQ. Be conversational, friendly, and helpful.
Learning style: {learning_style}
Instructions: {learning_instruction}

Personal learning preferences: {personalized_instruction}

Format your response properly using Markdown:
- Use **bold** for important concepts
- Use proper paragraph breaks for readability
- Use bullet points or numbered lists when appropriate
- Use headings with # symbols for section titles
- {model_instruction}

"""

        # Enhance prompt based on learning style
        if learning_style == "visual":
            prompt = generate_visual_prompt(base_prompt)
        else:
            prompt = base_prompt

        if context:
            prompt += f"""
Recent document content:
{context}

"""
        else:
            # No document, set expectations for basic chat
            prompt += """
The user has not uploaded any documents yet, so please respond to general questions.
If they ask about specific content, politely suggest they upload a document first.

"""

        if chat_context:
            prompt += f"""
Recent conversation:
{chat_context}

"""

        prompt += f"""
User's question: {user_message}

Please respond directly to the user's question.
"""

        if has_document:
            prompt += " If the question is about the document content, refer to it in your answer."
        
        prompt += """
Make your response well-structured and easy to read with proper formatting.
"""

        # Call Gemini model
        response = model.generate_content(prompt)
        ai_response = response.text
        
        # Post-process response based on learning style
        if learning_style == "visual":
            ai_response = add_visual_elements(ai_response)
        
        # Generate audio for auditory learners
        audio_url = None
        if learning_style == "auditory" and GTTS_AVAILABLE:
            try:
                logger.info(f"Generating audio for session {session_id}")
                audio_result = generate_audio_from_text(ai_response, session_id)
                if audio_result:
                    audio_url = audio_result["url"]
                    logger.info(f"Audio generated successfully: {audio_url}")
            except Exception as e:
                logger.error(f"Error generating audio: {str(e)}")
        
        # Add AI response to history
        session_data['chat_history'].append({
            'role': 'assistant',
            'content': ai_response,
            'timestamp': time.time(),
            'audio_url': audio_url
        })
        
        # Update last active time
        session_data['last_active'] = time.time()
        
        # Return the response with audio URL if available
        response_data = {
            "response": ai_response,
            "timestamp": time.time()
        }
        
        if audio_url:
            response_data["audio_url"] = audio_url
        
        return jsonify(response_data)
        
    except Exception as e:
        logger.error(f"Error in chat: {str(e)}")
        return jsonify({"error": "An error occurred. Please try again."}), 500

# Learning style API endpoint
@app.route('/api/learning-style', methods=['POST'])
def set_learning_style():
    """Set user's preferred learning style"""
    try:
        session_id = get_session_id()
        data = request.json
        
        if not data or 'style' not in data:
            return jsonify({"error": "No learning style provided"}), 400
            
        style = data['style']
        valid_styles = list(LEARNING_STYLES.keys())
        
        if style not in valid_styles:
            return jsonify({"error": f"Invalid learning style. Choose from: {', '.join(valid_styles)}"}), 400
        
        # Get previous style
        previous_style = user_sessions[session_id].get('learning_style', 'blended')
        
        # Update session
        user_sessions[session_id]['learning_style'] = style
        
        # If user is logged in, update their quiz data
        if 'user_email' in session:
            # Map learning style to quiz learningStyle field
            style_to_quiz_map = {
                "visual": "Watching videos",
                "reading": "Reading books",
                "auditory": "Listening to podcasts",
                "hands-on": "Doing it myself",
                "blended": "Others"
            }
            
            quiz_value = style_to_quiz_map.get(style, "Others")
            
            # Update quiz document
            quiz_collection.update_one(
                {'user_email': session['user_email']},
                {'$set': {'learningStyle': quiz_value}},
                upsert=True
            )
        
        # Add a message to the chat history if the style has changed
        if style != previous_style:
            style_info = LEARNING_STYLES.get(style, {})
            style_description = style_info.get("description", style.capitalize())
            
            message = f"Learning style changed to **{style}**. I'll now adapt my responses for {style_description}."
            
            # Add information about special features
            if style == "auditory" and GTTS_AVAILABLE:
                message += " I'll generate audio for my responses that you can play back."
            elif style == "visual":
                message += " I'll use visual organization, diagrams, and spatial cues in my explanations."
            
            user_sessions[session_id]['chat_history'].append({
                'role': 'assistant',
                'content': message,
                'timestamp': time.time()
            })
        
        return jsonify({
            "success": True,
            "learning_style": style,
            "tts_available": GTTS_AVAILABLE if style == "auditory" else None
        })
        
    except Exception as e:
        logger.error(f"Error in set_learning_style: {str(e)}")
        return jsonify({"error": "An error occurred. Please try again."}), 500

# Get documents API endpoint
@app.route('/api/documents', methods=['GET'])
def get_documents():
    """Get list of uploaded documents for current session"""
    try:
        session_id = get_session_id()
        
        # Return document list from session
        docs = user_sessions[session_id].get('documents', [])
        return jsonify({
            "documents": [
                {
                    "id": doc['id'],
                    "filename": doc['filename'],
                    "upload_time": doc['upload_time']
                } for doc in docs
            ]
        })
        
    except Exception as e:
        logger.error(f"Error in get_documents: {str(e)}")
        return jsonify({"error": "An error occurred. Please try again."}), 500

# Get chat history API endpoint
@app.route('/api/history', methods=['GET'])
def get_chat_history():
    """Get chat history for current session"""
    try:
        session_id = get_session_id()
        
        # Return chat history
        history = user_sessions[session_id].get('chat_history', [])
        return jsonify({
            "history": history
        })
        
    except Exception as e:
        logger.error(f"Error in get_chat_history: {str(e)}")
        return jsonify({"error": "An error occurred. Please try again."}), 500

# User's learning preferences API endpoint
@app.route('/api/user-preferences', methods=['GET'])
def get_user_preferences():
    """Get user's learning preferences from quiz results"""
    try:
        if 'user_email' not in session:
            return jsonify({"error": "Not logged in"}), 401
            
        quiz_data = quiz_collection.find_one({"user_email": session['user_email']})
        
        if not quiz_data:
            return jsonify({"error": "No quiz data found"}), 404
            
        # Remove MongoDB ID and email for security
        if '_id' in quiz_data:
            del quiz_data['_id']
        if 'user_email' in quiz_data:
            del quiz_data['user_email']
            
        return jsonify({
            "preferences": quiz_data
        })
        
    except Exception as e:
        logger.error(f"Error in get_user_preferences: {str(e)}")
        return jsonify({"error": "An error occurred. Please try again."}), 500

# Cleanup function for stale sessions
@app.before_request
def cleanup_sessions():
    """Clean up stale sessions (run periodically)"""
    now = time.time()
    # Only run cleanup every 5 minutes
    if now % 300 < 10:
        try:
            stale_threshold = now - (24 * 60 * 60)  # 24 hours
            stale_sessions = []
            
            # Identify stale sessions
            for session_id, session_data in user_sessions.items():
                if session_data['last_active'] < stale_threshold:
                    stale_sessions.append(session_id)
                    
            # Clean up stale sessions
            for session_id in stale_sessions:
                if session_id in user_sessions:
                    # Delete files
                    for doc in user_sessions[session_id].get('documents', []):
                        filepath = doc.get('filepath')
                        if filepath and os.path.exists(filepath):
                            os.remove(filepath)
                    
                    # Delete audio files
                    for entry in user_sessions[session_id].get('chat_history', []):
                        if 'audio_url' in entry and entry['audio_url']:
                            audio_file = entry['audio_url'].split('/')[-1]
                            audio_path = os.path.join(AUDIO_FOLDER, audio_file)
                            if os.path.exists(audio_path):
                                os.remove(audio_path)
                            
                    # Delete session data
                    del user_sessions[session_id]
                    
            if stale_sessions:
                logger.info(f"Cleaned up {len(stale_sessions)} stale sessions")
                
            # Also clean up old audio files (older than 6 hours)
            audio_threshold = now - (6 * 60 * 60)
            if os.path.exists(AUDIO_FOLDER):
                for filename in os.listdir(AUDIO_FOLDER):
                    filepath = os.path.join(AUDIO_FOLDER, filename)
                    if os.path.isfile(filepath) and os.path.getmtime(filepath) < audio_threshold:
                        os.remove(filepath)
                        logger.info(f"Cleaned up old audio file: {filepath}")
                
        except Exception as e:
            logger.error(f"Error in cleanup_sessions: {str(e)}")
# Add these new routes to app.py

# Teacher dashboard route
@app.route('/teacher')
def teacher_dashboard():
    if 'user_email' not in session or session.get('user_role') != 'teacher':
        flash('Please login as a teacher to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    return render_template('teacher_dashboard.html')

# Create quiz route
@app.route('/create_quiz', methods=['GET', 'POST'])
def create_quiz():
    if 'user_email' not in session or session.get('user_role') != 'teacher':
        flash('Please login as a teacher to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    if request.method == 'POST':
        # Process quiz creation form
        quiz_data = request.get_json()
        quiz_data['teacher_email'] = session['user_email']
        
        # Store quiz in database
        db.quizzes.insert_one(quiz_data)
        
        return jsonify({"success": True, "quiz_id": str(quiz_data.get('_id'))})
    
    return render_template('create_quiz.html')

# View student submissions route
@app.route('/student_submissions')
def student_submissions():
    if 'user_email' not in session or session.get('user_role') != 'teacher':
        flash('Please login as a teacher to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    # Fetch submissions that are for this teacher's quizzes
    teacher_quizzes = list(db.quizzes.find({"teacher_email": session['user_email']}))
    quiz_ids = [q['_id'] for q in teacher_quizzes]
    
    submissions = list(db.quiz_submissions.find({"quiz_id": {"$in": quiz_ids}}))
    
    return render_template('student_submissions.html', submissions=submissions, quizzes=teacher_quizzes)

# Create syllabus plan route
@app.route('/create_syllabus', methods=['GET', 'POST'])
def create_syllabus():
    if 'user_email' not in session or session.get('user_role') != 'teacher':
        flash('Please login as a teacher to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    if request.method == 'POST':
        syllabus_data = request.get_json()
        syllabus_data['teacher_email'] = session['user_email']
        
        # Store syllabus in database
        db.syllabi.insert_one(syllabus_data)
        
        return jsonify({"success": True})
    
    return render_template('create_syllabus.html')


@app.route('/admin')
def admin_dashboard():
    if 'user_email' not in session or session.get('user_role') != 'admin':
        flash('Please login as admin to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    # Get all users for admin view
    all_users = list(users_collection.find({}, {'password': 0}))  # Exclude password field
    
    return render_template('admin_dashboard.html', users=all_users)

# Admin edit user route
@app.route('/admin/edit_user/<user_email>', methods=['GET', 'POST'])
def admin_edit_user(user_email):
    if 'user_email' not in session or session.get('user_role') != 'admin':
        flash('Please login as admin to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    if request.method == 'POST':
        # Process form data
        new_email = request.form.get('email')
        new_name = request.form.get('full_name')
        new_role = request.form.get('role')
        
        # Update user data
        update_data = {
            "email": new_email,
            "full_name": new_name,
            "role": new_role
        }
        
        # Check if password is being changed
        new_password = request.form.get('password')
        if new_password and len(new_password) >= 6:
            update_data["password"] = generate_password_hash(new_password)
        
        # Update user in database
        users_collection.update_one(
            {"email": user_email},
            {"$set": update_data}
        )
        
        # If email is being changed, we need to update related records
        if new_email != user_email:
            # Update user email in quiz results
            quiz_collection.update_many(
                {"user_email": user_email},
                {"$set": {"user_email": new_email}}
            )
            
            # Update any other collections that reference user email
            
        flash('User information updated successfully.', 'success')
        return redirect(url_for('admin_dashboard'))
    
    # Get user data for edit form
    user = users_collection.find_one({"email": user_email})
    if not user:
        flash('User not found.', 'error')
        return redirect(url_for('admin_dashboard'))
    
    return render_template('admin_edit_user.html', user=user)

# Admin create user route
@app.route('/admin/create_user', methods=['GET', 'POST'])
def admin_create_user():
    if 'user_email' not in session or session.get('user_role') != 'admin':
        flash('Please login as admin to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    if request.method == 'POST':
        # Process form data
        full_name = request.form.get('full_name')
        email = request.form.get('email')
        password = request.form.get('password')
        role = request.form.get('role')
        
        # Validation
        if not all([full_name, email, password, role]):
            flash('All fields are required.', 'error')
            return redirect(url_for('admin_create_user'))
        
        if not is_valid_email(email):
            flash('Please enter a valid email address.', 'error')
            return redirect(url_for('admin_create_user'))
        
        if len(password) < 6:
            flash('Password must be at least 6 characters long.', 'error')
            return redirect(url_for('admin_create_user'))
        
        # Check if user already exists
        if users_collection.find_one({"email": email}):
            flash('Email already exists. Please use a different email.', 'error')
            return redirect(url_for('admin_create_user'))
        
        # Hash the password
        hashed_password = generate_password_hash(password)
        
        # Create new user
        new_user = {
            "full_name": full_name,
            "email": email,
            "password": hashed_password,
            "role": role
        }
        
        # Insert into database
        users_collection.insert_one(new_user)
        
        flash('User created successfully.', 'success')
        return redirect(url_for('admin_dashboard'))
    
    return render_template('admin_create_user.html')

# Admin delete user route
@app.route('/admin/delete_user/<user_email>', methods=['POST'])
def admin_delete_user(user_email):
    if 'user_email' not in session or session.get('user_role') != 'admin':
        flash('Please login as admin to access this page.', 'error')
        return redirect(url_for('loginpage'))
    
    # Prevent admin from deleting themselves
    if user_email == session['user_email']:
        flash('You cannot delete your own admin account.', 'error')
        return redirect(url_for('admin_dashboard'))
    
    # Delete user from database
    users_collection.delete_one({"email": user_email})
    
    # Delete related user data
    quiz_collection.delete_many({"user_email": user_email})
    
    flash('User deleted successfully.', 'success')
    return redirect(url_for('admin_dashboard'))

# Quiz generation API
@app.route('/api/generate_quiz', methods=['POST'])
def generate_quiz_api():
    if 'user_email' not in session:
        return jsonify({"error": "You must be logged in to use this feature"}), 401
        
    try:
        data = request.get_json()
        document_id = data.get('document_id')
        learning_style = data.get('learning_style', 'blended')
        question_count = data.get('question_count', 5)
        
        # Fetch the document content
        doc_content = document_content.get(document_id, '')
        
        if not doc_content:
            return jsonify({"error": "Document not found or empty"}), 404
        
        # Get user's quiz results for personalization
        quiz_results = quiz_collection.find_one({"user_email": session['user_email']})
        
        # Build prompt for Gemini API based on learning style and preferences
        prompt_parts = [
            f"Generate {question_count} multiple-choice quiz questions about the following content.",
            "Each question should have exactly 4 options with one correct answer.",
            f"Format the response as a JSON array of objects with the structure: [{{'question': 'Question text', 'options': ['option1', 'option2', 'option3', 'option4'], 'correct_answer': 0}}] where correct_answer is the index (0-3) of the correct option.",
        ]
        
        # Add learning style specific instructions
        if learning_style == 'visual':
            prompt_parts.append("Make the questions focused on relationships, patterns, and visual concepts.")
        elif learning_style == 'auditory':
            prompt_parts.append("Phrase questions conversationally, focusing on dialogue and verbal concepts.")
        elif learning_style == 'hands-on':
            prompt_parts.append("Focus questions on practical applications, problem-solving, and step-by-step processes.")
        elif learning_style == 'reading':
            prompt_parts.append("Create detailed, text-focused questions that test comprehension and analytical reading skills.")
        
        # Add personalization based on quiz results
        if quiz_results:
            interaction_style = quiz_results.get('chatbotInteraction', '')
            if interaction_style == 'Witty and humorous':
                prompt_parts.append("Add a light touch of humor to the questions where appropriate.")
            
            struggle_help = quiz_results.get('struggleHelp', '')
            if struggle_help == 'Hints or clues':
                prompt_parts.append("Include subtle hints within the options for challenging questions.")
            
            focus_duration = quiz_results.get('studyDuration', '')
            if focus_duration == 'less than 15 minutes':
                prompt_parts.append("Keep questions concise and straightforward.")
            elif focus_duration == 'More than an hour':
                prompt_parts.append("Include some more complex, multi-step thinking questions.")
        
        # Combine prompt parts
        prompt = "\n".join(prompt_parts) + "\n\nContent:\n" + doc_content[:5000]  # Limit content length
        
        # Call Gemini API
        if HAS_GEMINI:
            response = model.generate_content(prompt)
            response_text = response.text
            
            # Extract JSON from response
            try:
                # Find JSON array in the response
                json_match = re.search(r'\[\s*\{.*\}\s*\]', response_text, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                    questions = json.loads(json_str)
                    
                    return jsonify({
                        "success": True,
                        "questions": questions
                    })
                else:
                    return jsonify({"error": "Failed to parse generated questions"}), 500
            except Exception as e:
                logger.error(f"Error parsing quiz questions: {str(e)}")
                return jsonify({"error": "Failed to parse generated questions"}), 500
        else:
            return jsonify({"error": "AI model is not available"}), 503
            
    except Exception as e:
        logger.error(f"Error generating quiz: {str(e)}")
        return jsonify({"error": "An error occurred while generating the quiz"}), 500

# Main entry point
if __name__ == '__main__':
    # Ensure directories exist
    for directory in ['uploads', 'static', os.path.join('static', 'audio'), 'templates']:
        if not os.path.exists(directory):
            os.makedirs(directory)
    
    # Run the app
    port = int(os.environ.get("PORT", 8000))
    debug_mode = os.environ.get("FLASK_ENV") != "production"
    app.run(debug=debug_mode, host='0.0.0.0', port=port)
